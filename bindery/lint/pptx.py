"""PPTX token-compliance linting — M3-spec.md §2.2."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from bindery.ds.loader import DesignSystem


def lint(artifact_path: Path, ds: DesignSystem) -> list:
    from bindery.lint import LintViolation, allowed_values

    allowed = allowed_values(ds)
    violations: list[LintViolation] = []

    prs = Presentation(artifact_path)
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                # Check both paragraph-level font (where this codebase's own
                # components set color/size, e.g. `p.font.color.rgb = ...` in
                # design-systems/*/components/pptx/*.py) and each run's font
                # (in case a future component sets run-level formatting
                # directly) — checking only one level silently misses the
                # other, which is exactly how this went unnoticed at first.
                fonts = [(f"slide {slide_idx} / shape {shape_idx} / paragraph {para_idx}", paragraph.font)]
                for run_idx, run in enumerate(paragraph.runs):
                    fonts.append(
                        (f"slide {slide_idx} / shape {shape_idx} / run {run_idx}", run.font)
                    )

                for location, font in fonts:
                    if font.color and font.color.type is not None:
                        try:
                            hexval = f"#{font.color.rgb}"
                        except (AttributeError, TypeError):
                            hexval = None
                        if hexval and hexval not in allowed:
                            violations.append(LintViolation(location, "color", hexval))

                    if font.size is not None:
                        size_str = str(int(font.size.pt))
                        if size_str not in allowed:
                            violations.append(LintViolation(location, "font-size", size_str))

                    if font.name is not None and font.name not in allowed:
                        violations.append(LintViolation(location, "font-family", font.name))

    return violations
