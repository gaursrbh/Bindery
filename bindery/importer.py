"""DS importer — mainPRD R13, M4-spec.md §4, extended by issue #43.

Extracts a *candidate* tokens.json for human review. Never installs a live
design system — components and rules still need human authoring; this is
token intake, not authoring (M4-spec.md §4.2). Two input paths: an existing
PPTX deck (M4, network-free) and a live website's CSS (issue #43 — the one
input path in this module that touches the network; the importer is a
one-off intake/dev-tool step, not part of the deterministic render
pipeline mainPRD §6.4's "no network" purity guarantee governs).
"""

from __future__ import annotations

import json
import re
import urllib.request
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

_HEX_RE = re.compile(r"#[0-9a-fA-F]{6}\b|#[0-9a-fA-F]{3}\b")
_FONT_FAMILY_RE = re.compile(r"font-family\s*:\s*([^;\"'}]+)")
_FONT_SIZE_RE = re.compile(r"font-size\s*:\s*(\d+)px")
# CSS keywords and custom-property references aren't real font names — a
# real site (anthropic.com) surfaced both "inherit" and a raw
# "var(--...)" reference as if they were fonts before this filter existed.
_FONT_KEYWORD_BLOCKLIST = {"inherit", "initial", "unset", "none", "revert"}


@dataclass
class ImportReport:
    colors: Counter = field(default_factory=Counter)
    sizes: Counter = field(default_factory=Counter)
    fonts: Counter = field(default_factory=Counter)


def scan_pptx(path: Path) -> ImportReport:
    report = ImportReport()
    prs = Presentation(path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                fonts_to_check = [paragraph.font] + [r.font for r in paragraph.runs]
                for font in fonts_to_check:
                    if font.color and font.color.type is not None:
                        try:
                            report.colors[f"#{font.color.rgb}"] += 1
                        except (AttributeError, TypeError):
                            pass
                    if font.size is not None:
                        report.sizes[int(font.size.pt)] += 1
                    if font.name is not None:
                        report.fonts[font.name] += 1

    return report


def scan_website(url: str, timeout: int = 15) -> ImportReport:
    """Fetches the page's raw HTML and regex-scans it for hex colors and
    font declarations in inline <style> blocks and style="" attributes.

    Deliberate scope cut: does not fetch externally linked stylesheets
    (<link rel="stylesheet">) — a real second HTTP round-trip per
    stylesheet, and most marketing sites inline their critical CSS or use a
    component framework whose color choices show up in inline styles
    anyway. A page whose entire palette lives in an external .css file will
    under-report; that's a known gap, not a silent one.
    """
    report = ImportReport()
    req = urllib.request.Request(url, headers={"User-Agent": "Bindery-DS-Importer/0.1"})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        html = resp.read().decode(resp.headers.get_content_charset() or "utf-8", errors="replace")

    for m in _HEX_RE.finditer(html):
        report.colors[m.group().upper()] += 1
    for m in _FONT_FAMILY_RE.finditer(html):
        # First comma-separated family in the stack, quotes stripped.
        first = m.group(1).split(",")[0].strip().strip("'\"")
        if first and first.lower() not in _FONT_KEYWORD_BLOCKLIST and not first.startswith("var("):
            report.fonts[first] += 1
    for m in _FONT_SIZE_RE.finditer(html):
        report.sizes[int(m.group(1))] += 1

    return report


def candidate_tokens(report: ImportReport) -> dict:
    ranked_colors = [c for c, _ in report.colors.most_common()]
    ranked_sizes = sorted(report.sizes.most_common(), key=lambda kv: -kv[1])
    top_font = report.fonts.most_common(1)[0][0] if report.fonts else "Helvetica Neue"

    color_slots = ["primary", "secondary", "neutral", "background", "text"]
    color_tokens = {
        slot: {"value": ranked_colors[i]}
        for i, slot in enumerate(color_slots)
        if i < len(ranked_colors)
    }

    size_names = ["headline-size", "stat-value-size", "eyebrow-size", "stat-label-size"]
    typography_tokens = {"family": {"value": top_font}}
    for i, name in enumerate(size_names):
        if i < len(ranked_sizes):
            typography_tokens[name] = {"value": str(ranked_sizes[i][0])}

    return {"color": color_tokens, "typography": typography_tokens}


def write_candidate(deck_path: Path, out_dir: Path) -> Path:
    report = scan_pptx(deck_path)
    tokens = candidate_tokens(report)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "candidate-tokens.json"
    out_path.write_text(json.dumps(tokens, indent=2))
    return out_path


def write_candidate_from_website(url: str, out_dir: Path) -> Path:
    report = scan_website(url)
    tokens = candidate_tokens(report)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "candidate-tokens.json"
    out_path.write_text(json.dumps(tokens, indent=2))
    return out_path
