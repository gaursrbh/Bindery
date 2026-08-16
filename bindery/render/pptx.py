"""(Composition, DesignSystem) -> Artifact — mainPRD R4, §6.4, M0-spec.md §4.

Pure and deterministic per §6.4: no filesystem/network/wall-clock reads
beyond the explicit `out_path` write, no randomness.
"""

from __future__ import annotations

import json
import time
from dataclasses import dataclass
from pathlib import Path

from jsonschema import Draft202012Validator
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from referencing import Registry, Resource

from bindery.ds.loader import SCHEMA_ROOT, DesignSystem
from bindery.render.errors import CompositionError, RenderError
from bindery.render.overflow import check_overflow

_SLIDE_WIDTH_IN = 10
_SLIDE_HEIGHT_IN = 5.63
_BLANK_LAYOUT_INDEX = 6
_MARGIN_TOP_IN = 0.5

_CORE_SCHEMA = json.loads((SCHEMA_ROOT / "core.schema.json").read_text())
_REGISTRY = Registry().with_resource(
    "core.schema.json", Resource.from_contents(_CORE_SCHEMA)
)


@dataclass
class RenderResult:
    path: Path
    duration_ms: int
    blocks_rendered: int


def _hexcolor(tokens: dict, name: str) -> RGBColor:
    return RGBColor.from_string(tokens["color"][name]["value"].lstrip("#"))


def _validate(composition: dict, ds: DesignSystem) -> dict:
    target = composition.get("target")
    schema = ds.effective_schemas.get(target)
    if schema is None:
        raise CompositionError(
            f"design system '{ds.spec}' has no schema for target {target!r}; "
            f"available targets: {sorted(ds.effective_schemas)}"
        )

    # Strip $id so ref resolution starts from an empty base URI, matching the
    # "core.schema.json" key the registry resource is stored under — schemas
    # here use filename-relative $refs, not real URIs.
    schema_for_validation = {k: v for k, v in schema.items() if k != "$id"}
    validator = Draft202012Validator(schema_for_validation, registry=_REGISTRY)
    errors = sorted(validator.iter_errors(composition), key=lambda e: list(e.path))
    if errors:
        first = errors[0]
        location = "/".join(str(p) for p in first.path) or "<root>"
        available = sorted(ds.layout_fns.get(target, {}))
        raise CompositionError(
            f"composition invalid at {location}: {first.message}; "
            f"available components: {available}"
        )
    return schema


def render(composition: dict, ds: DesignSystem, out_path: Path) -> RenderResult:
    start = time.monotonic()
    _validate(composition, ds)

    target = composition["target"]
    layout_fns = ds.layout_fns[target]
    tokens = ds.tokens
    font_family = tokens.get("typography", {}).get("family", {}).get("value", "Helvetica")

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

    def _new_slide():
        s = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT_INDEX])
        bg = s.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hexcolor(tokens, "background")
        return s

    slide = None
    cursor_y = _MARGIN_TOP_IN
    for index, block in enumerate(composition["blocks"]):
        component = block["component"]
        layout_fn = layout_fns[component]

        # A `title` block starts a new slide — mirrors how a real deck is
        # authored (each title is a new slide's headline), and is what
        # M0/M1's fixed-position renderer never modeled: it placed every
        # block on one slide at hardcoded coordinates, so 2+ blocks of the
        # same component type landed on identical coordinates and rendered
        # as garbled overlapping text (found via a real brief with 8 blocks
        # across 3 titles — confirmed visually with LibreOffice rasterization,
        # not by the test suite, which only ever exercised 1-2 blocks).
        if component == "title" or slide is None:
            slide = _new_slide()
            cursor_y = _MARGIN_TOP_IN

        shapes_before = len(slide.shapes)
        consumed = layout_fn(slide, block["props"], tokens, cursor_y)
        cursor_y += consumed if consumed is not None else 0.0
        for shape in list(slide.shapes)[shapes_before:]:
            if shape.has_text_frame:
                check_overflow(shape, font_family, index, component)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    return RenderResult(
        path=out_path,
        duration_ms=int((time.monotonic() - start) * 1000),
        blocks_rendered=len(composition["blocks"]),
    )
