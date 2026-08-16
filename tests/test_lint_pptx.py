from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from bindery.ds import loader
from bindery.lint.pptx import lint
from bindery.render.pptx import render


def test_compliant_render_has_no_violations(ds_root, tmp_path):
    ds = loader.load("reference", root=ds_root)
    composition = {
        "schema": "bindery/v1", "design_system": "reference@1.0.0", "target": "pptx",
        "blocks": [{"component": "title", "props": {"headline": "Q3 update", "eyebrow": "fy26"}}],
    }
    out = tmp_path / "out.pptx"
    render(composition, ds, out)
    assert lint(out, ds) == []


def test_off_token_color_and_size_detected(ds_root, tmp_path):
    ds = loader.load("reference", root=ds_root)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "hi"
    p.font.color.rgb = RGBColor.from_string("FF00FF")
    p.font.size = Pt(99)
    out = tmp_path / "bad.pptx"
    prs.save(out)

    violations = lint(out, ds)
    kinds = {v.kind for v in violations}
    assert "color" in kinds
    assert "font-size" in kinds
