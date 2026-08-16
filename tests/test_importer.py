import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from bindery.importer import candidate_tokens, scan_pptx, write_candidate


def _make_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "Headline"
    p.font.color.rgb = RGBColor.from_string("1F3A5F")
    p.font.size = Pt(32)

    box2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "Body"
    p2.font.color.rgb = RGBColor.from_string("1F3A5F")
    p2.font.size = Pt(12)

    prs.save(path)


def test_scan_pptx_counts_colors_and_sizes(tmp_path):
    deck = tmp_path / "deck.pptx"
    _make_deck(deck)

    report = scan_pptx(deck)
    assert report.colors["#1F3A5F"] == 2
    assert report.sizes[32] == 1
    assert report.sizes[12] == 1


def test_candidate_tokens_shape(tmp_path):
    deck = tmp_path / "deck.pptx"
    _make_deck(deck)
    report = scan_pptx(deck)

    tokens = candidate_tokens(report)
    assert tokens["color"]["primary"]["value"] == "#1F3A5F"
    assert "family" in tokens["typography"]


def test_write_candidate_writes_file(tmp_path):
    deck = tmp_path / "deck.pptx"
    _make_deck(deck)
    out_dir = tmp_path / "out"

    out_path = write_candidate(deck, out_dir)
    assert out_path.exists()
    data = json.loads(out_path.read_text())
    assert "color" in data and "typography" in data
