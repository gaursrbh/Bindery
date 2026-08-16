import pytest
from pptx import Presentation

from bindery.ds import loader
from bindery.render.errors import CompositionError, RenderError
from bindery.render.pptx import render


@pytest.fixture
def ds(ds_root):
    return loader.load("reference", root=ds_root)


def _composition(*blocks):
    return {
        "schema": "bindery/v1",
        "design_system": "reference@1.0.0",
        "target": "pptx",
        "blocks": list(blocks),
    }


def test_render_title(ds, tmp_path):
    composition = _composition(
        {"component": "title", "props": {"headline": "Q3 board update", "eyebrow": "quarterly"}}
    )
    result = render(composition, ds, tmp_path / "out.pptx")
    assert result.path.exists()
    assert result.blocks_rendered == 1
    prs = Presentation(result.path)
    assert len(prs.slides) == 1


def test_render_stat_trio(ds, tmp_path):
    composition = _composition(
        {
            "component": "stat-trio",
            "props": {
                "stats": [
                    {"value": "$4.2M", "label": "Revenue", "delta": "+12%"},
                    {"value": "312", "label": "New logos"},
                    {"value": "94%", "label": "Retention"},
                ]
            },
        }
    )
    result = render(composition, ds, tmp_path / "out.pptx")
    assert result.path.exists()


def test_render_bullet_list(ds, tmp_path):
    composition = _composition(
        {
            "component": "bullet-list",
            "props": {"heading": "Highlights", "items": ["Shipped M0", "Closed 3 decisions"]},
        }
    )
    result = render(composition, ds, tmp_path / "out.pptx")
    assert result.path.exists()


def test_render_image_callout(ds, tmp_path):
    composition = _composition(
        {
            "component": "image-callout",
            "props": {"asset": "chart.png", "caption": "Growth by segment"},
        }
    )
    result = render(composition, ds, tmp_path / "out.pptx")
    assert result.path.exists()


def test_render_multi_block(ds, tmp_path):
    composition = _composition(
        {"component": "title", "props": {"headline": "Q3 board update"}},
        {
            "component": "stat-trio",
            "props": {
                "stats": [
                    {"value": "$4.2M", "label": "Revenue"},
                    {"value": "312", "label": "New logos"},
                    {"value": "94%", "label": "Retention"},
                ]
            },
        },
    )
    result = render(composition, ds, tmp_path / "out.pptx")
    assert result.blocks_rendered == 2


def test_invalid_composition_raises_with_available_components(ds, tmp_path):
    composition = _composition({"component": "nonexistent", "props": {}})
    with pytest.raises(CompositionError) as exc:
        render(composition, ds, tmp_path / "out.pptx")
    message = str(exc.value)
    assert "title" in message and "stat-trio" in message


def test_overflow_raises_render_error(ds, tmp_path):
    composition = _composition(
        {
            "component": "bullet-list",
            "props": {
                "heading": "This heading is long enough to wrap onto two lines of text",
                "items": [
                    "A" + " word" * 23,
                    "B" + " word" * 23,
                    "C" + " word" * 23,
                    "D" + " word" * 23,
                    "E" + " word" * 23,
                    "F" + " word" * 23,
                ],
            },
        }
    )
    with pytest.raises(RenderError) as exc:
        render(composition, ds, tmp_path / "out.pptx")
    assert exc.value.block_index == 0
    assert exc.value.prop == "bullet-list"


def test_title_starts_new_slide_and_blocks_stack(ds, tmp_path):
    """Regression test: M0/M1's renderer placed every block on one slide at
    hardcoded per-component-type coordinates. A composition with more than
    one title (or more than one of the same component) rendered every
    instance of a component on top of the others — found via a real 8-block
    brief, confirmed by rasterizing with LibreOffice."""
    composition = _composition(
        {"component": "title", "props": {"headline": "Slide one"}},
        {"component": "bullet-list", "props": {"items": ["a", "b"]}},
        {"component": "title", "props": {"headline": "Slide two"}},
        {"component": "bullet-list", "props": {"items": ["c", "d"]}},
        {"component": "bullet-list", "props": {"items": ["e", "f"]}},
    )
    result = render(composition, ds, tmp_path / "out.pptx")
    assert result.blocks_rendered == 5

    prs = Presentation(result.path)
    assert len(prs.slides) == 2  # one per title block

    slide2_shapes = list(prs.slides[1].shapes)
    assert len(slide2_shapes) == 3  # title + 2 bullet-lists, not overlapping

    # No two shapes on the same slide start at the same y — the exact
    # symptom of the bug (every bullet-list landed at the same fixed y).
    tops = [s.top for s in slide2_shapes]
    assert len(tops) == len(set(tops))
