"""(Composition, tokens.json) -> .pptx, per mainPRD §6.4.

Only handles the two shared components (title, stat-trio) — enough to test
whether the shared bindery/v1 core survives an actual render, not to be a
real renderer.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).parent


def hexcolor(tokens, name):
    return RGBColor.from_string(tokens["color"][name]["value"].lstrip("#"))


def render_title(slide, props, tokens):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    if "eyebrow" in props:
        p = tf.paragraphs[0]
        p.text = props["eyebrow"].upper()
        p.font.size = Pt(int(tokens["typography"]["eyebrow-size"]["value"]))
        p.font.color.rgb = hexcolor(tokens, props.get("accent", "primary"))
        p.font.bold = True
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    p2.text = props["headline"]
    p2.font.size = Pt(int(tokens["typography"]["headline-size"]["value"]))
    p2.font.color.rgb = hexcolor(tokens, "text")
    p2.font.bold = True


def render_stat_trio(slide, props, tokens):
    stats = props["stats"]
    col_width = Inches(3.0)
    for i, stat in enumerate(stats):
        left = Inches(0.6) + i * col_width
        box = slide.shapes.add_textbox(left, Inches(2.6), col_width - Inches(0.2), Inches(2))
        tf = box.text_frame
        tf.word_wrap = True
        p_val = tf.paragraphs[0]
        p_val.text = stat["value"]
        p_val.font.size = Pt(int(tokens["typography"]["stat-value-size"]["value"]))
        p_val.font.bold = True
        p_val.font.color.rgb = hexcolor(tokens, "primary")

        p_label = tf.add_paragraph()
        p_label.text = stat["label"]
        p_label.font.size = Pt(int(tokens["typography"]["stat-label-size"]["value"]))
        p_label.font.color.rgb = hexcolor(tokens, "neutral")

        if "delta" in stat:
            p_delta = tf.add_paragraph()
            p_delta.text = stat["delta"]
            p_delta.font.size = Pt(int(tokens["typography"]["stat-label-size"]["value"]))
            p_delta.font.color.rgb = hexcolor(tokens, "secondary")


RENDERERS = {"title": render_title, "stat-trio": render_stat_trio}


def main():
    comp_path = sys.argv[1] if len(sys.argv) > 1 else "composition-pptx.json"
    out_path = sys.argv[2] if len(sys.argv) > 2 else "out-pptx.pptx"
    composition = json.loads((HERE / comp_path).read_text())
    tokens = json.loads((HERE / "tokens.json").read_text())
    assert composition["target"] == "pptx", "wrong target for pptx renderer"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = hexcolor(tokens, "background")

    for block in composition["blocks"]:
        renderer = RENDERERS.get(block["component"])
        if renderer is None:
            raise ValueError(f"pptx renderer has no component {block['component']!r}")
        renderer(slide, block["props"], tokens)

    prs.save(HERE / out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
